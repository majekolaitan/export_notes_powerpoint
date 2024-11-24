import os
from pptx import Presentation
from docx import Document

def extract_notes_to_docx(input_dir, output_dir):
    """
    Extracts notes from all PowerPoint presentations in a directory and saves them to Word documents.
    
    :param input_dir: Directory containing PowerPoint files.
    :param output_dir: Directory to save Word documents with extracted notes.
    """
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    for filename in os.listdir(input_dir):
        if filename.endswith(".pptx"):
            ppt_path = os.path.join(input_dir, filename)
            docx_path = os.path.join(output_dir, f"{os.path.splitext(filename)[0]}_notes.docx")

            # Load PowerPoint presentation
            presentation = Presentation(ppt_path)

            # Create a Word document
            document = Document()
            document.add_heading(f"Notes from {filename}", level=1)

            for i, slide in enumerate(presentation.slides):
                # Add slide title to Word document
                slide_title = slide.shapes.title.text if slide.shapes.title else f"Slide {i+1}"
                document.add_heading(slide_title, level=2)

                # Extract and add slide notes to Word document
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text
                    document.add_paragraph(notes_text)
                else:
                    document.add_paragraph("(No notes)")

            # Save Word document
            document.save(docx_path)
            print(f"Extracted notes from {filename} to {docx_path}")

# Define input and output directories
input_directory = "powerpoints"  # Replace with the path to your PowerPoint files
output_directory = "documents"  # Replace with the path to save Word documents

# Run the extraction
extract_notes_to_docx(input_directory, output_directory)
